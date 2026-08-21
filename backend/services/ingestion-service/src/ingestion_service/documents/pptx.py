from __future__ import annotations

import io
from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from ingestion_service.domain import Asset, BlockType, DocumentUnit, ExtractedBlock


class PPTXReader:
    def supports(self, path: Path, content_type: str) -> bool:
        return path.suffix.lower() == ".pptx" or content_type == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    def units(self, path: Path) -> Iterator[DocumentUnit]:
        presentation = Presentation(path)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            blocks: list[ExtractedBlock] = []
            assets: list[Asset] = []
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
                if title:
                    blocks.append(
                        ExtractedBlock(
                            block_type=BlockType.HEADING,
                            text=title,
                            section_path=[title],
                        )
                    )
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if getattr(shape, "has_text_frame", False):
                    text = "\n".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    if text:
                        blocks.append(
                            ExtractedBlock(
                                block_type=BlockType.TEXT,
                                text=text,
                                section_path=[title] if title else [],
                            )
                        )
                if getattr(shape, "has_table", False):
                    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    blocks.append(
                        ExtractedBlock(
                            block_type=BlockType.TABLE,
                            text="\n".join(" | ".join(row) for row in rows),
                            section_path=[title] if title else [],
                        )
                    )
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    data = shape.image.blob
                    with Image.open(io.BytesIO(data)) as image:
                        width, height = image.size
                        extension = (image.format or shape.image.ext or "png").lower()
                    assets.append(
                        Asset(
                            data=data,
                            extension=extension,
                            content_type=shape.image.content_type,
                            width=width,
                            height=height,
                        )
                    )
            yield DocumentUnit(
                number=slide_number,
                label=f"slide {slide_number}",
                blocks=blocks,
                assets=assets,
                metadata={"format": "pptx"},
            )
